import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import pgeocode
import pydeck as pdk
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans, DBSCAN
from sklearn.linear_model import LinearRegression
import json
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile

st.set_page_config(layout="wide")
st.title("Smile Agent (POC) – Rapport BI & Agent Conversationnel")

# --- 1. Upload des 3 fichiers + modèle PPTX ---
st.sidebar.header("Chargement des fichiers")
file_tx = st.sidebar.file_uploader("Transactions (Test)", ["csv","xlsx"], key="tx")
file_merch = st.sidebar.file_uploader("Caractéristiques marchands", ["csv","xlsx"], key="merch")
file_weather = st.sidebar.file_uploader("Données météo", ["csv","xlsx"], key="weather")
file_template = st.sidebar.file_uploader("Modèle PPTX (exemple.pptx)", ["pptx"], key="template")

@st.cache_data
def read_file(u):
    if u is None:
        return None
    if u.name.lower().endswith('.csv'):
        try:
            return pd.read_csv(u, sep=None, engine='python')
        except:
            return pd.read_csv(u)
    return pd.read_excel(u, engine='openpyxl')

# Lecture des données
df_tx = read_file(file_tx)
df_merch = read_file(file_merch)
df_weather = read_file(file_weather)

if df_tx is not None and df_merch is not None and df_weather is not None:
    st.success("✅ Tous les fichiers chargés")

    # --- 2. Préparation et fusion ---
    df_tx['MONTANT'] = (
        df_tx['MONTANT'].astype(str)
               .str.replace(r"[^0-9,.-]", '', regex=True)
               .str.replace(',', '.')
               .astype(float)
    )
    df_tx['DATETIME'] = pd.to_datetime(
        df_tx['DATE'].astype(str) + ' ' + df_tx['HEURE'].astype(str),
        dayfirst=True, errors='coerce'
    )
    df_tx['HOUR'] = df_tx['DATETIME'].dt.hour
    df_tx['DAY'] = df_tx['DATETIME'].dt.date
    df_tx['WEEK'] = df_tx['DATETIME'].dt.to_period('W').apply(lambda r: r.start_time)
    df_tx['MONTH'] = df_tx['DATETIME'].dt.to_period('M').apply(lambda r: r.start_time)

    df_weather.columns = df_weather.columns.str.strip().str.replace(r"\s+", ' ', regex=True).str.upper()
    if 'TEMPÉRATURE' in df_weather.columns:
        df_weather.rename(columns={'TEMPÉRATURE':'TEMP'}, inplace=True)
    if 'CODE_POSTAL' in df_weather.columns and 'CODE POSTAL' not in df_weather.columns:
        df_weather.rename(columns={'CODE_POSTAL':'CODE POSTAL'}, inplace=True)

    df = (
        df_tx
        .merge(
            df_merch.rename(columns={'Organization_type':'TYPE_COMMERCE'})[['REF_MARCHAND','TYPE_COMMERCE']],
            on='REF_MARCHAND', how='left'
        )
        .merge(df_weather[['CODE POSTAL','TEMP']], on='CODE POSTAL', how='left')
    )

    nomi = pgeocode.Nominatim('fr')
    unique_cp = df['CODE POSTAL'].astype(str).str.zfill(5).drop_duplicates().tolist()
    geo = nomi.query_postal_code(unique_cp)
    df_geo = geo[['postal_code','latitude','longitude']].dropna()
    df_geo['postal_code'] = df_geo['postal_code'].astype(str).str.zfill(5)
    df_geo.rename(columns={'postal_code':'CODE POSTAL'}, inplace=True)

    # --- 3. Rapport BI détaillé et génération PPTX ---
    if st.button("📄 Smile Magic Report"):
        # 1. Indicateurs descriptifs
        st.header("1. Indicateurs descriptifs")
        total_tx = len(df)
        by_day = df.groupby('DAY').size()
        by_week = df.groupby('WEEK').size()
        by_month = df.groupby('MONTH').size()
        ca_total = df['MONTANT'].sum()
        panier_moy = ca_total / total_tx
        st.metric("Transactions totales", f"{total_tx:,}")
        st.metric("Chiffre d'affaires total (€)", f"{ca_total:,.2f}")
        st.metric("Panier moyen (€)", f"{panier_moy:,.2f}")

        # Périodicité
        if by_day.size == 1:
            st.subheader(f"Analyse journalière ({by_day.index[0]})")
            st.write(int(by_day.iloc[0]))
        else:
            cols = st.columns(3)
            cols[0].write(by_day.to_frame('Nb TX'))
            cols[1].write(by_week.to_frame('Nb TX'))
            cols[2].write(by_month.to_frame('Nb TX'))
            if by_month.size > 1:
                st.subheader("Évolution mensuelle (T/T-1)")
                st.line_chart(by_month.pct_change().fillna(0))

        # Distribution montants
        st.subheader("Distribution des montants")
        stats = df['MONTANT'].describe(percentiles=[0.1,0.25,0.5,0.75,0.9])
        st.table(stats[['10%','25%','50%','75%','90%','mean']]
                 .rename({'10%':'P10','25%':'P25','50%':'Médiane','75%':'P75','90%':'P90','mean':'Moyenne'}))

        # Répartition par type et top-performers
        st.subheader("Répartition par type de commerce et top-performers")
        ca_type = df.groupby('TYPE_COMMERCE')['MONTANT'].agg(['sum','count'])
        ca_type['panier_moy'] = ca_type['sum']/ca_type['count']
        st.dataframe(ca_type)
        top_tx = df['MARCHAND'].value_counts().head(5)
        top_ca = df.groupby('MARCHAND')['MONTANT'].sum().nlargest(5)
        st.write("**Top 5 par volume:**", top_tx)
        st.write("**Top 5 par CA:**", top_ca)

        # Capture figure descriptifs
        fig_desc, ax_desc = plt.subplots()
        ca_type['count'].plot(kind='bar', ax=ax_desc, title='Nb TX par type de commerce')
        plt.tight_layout(); st.pyplot(fig_desc)

        # 2. Temporalité fine
        st.header("2. Temporalité fine")
        hourly = df.groupby('HOUR').agg(nb=('MONTANT','size'), avg=('MONTANT','mean'))
        fig_temporal_nb, ax_nb = plt.subplots()
        hourly['nb'].plot(kind='bar', ax=ax_nb, title='Transactions par heure')
        plt.tight_layout(); st.pyplot(fig_temporal_nb)
        fig_temporal_avg, ax_avg = plt.subplots()
        hourly['avg'].plot(kind='bar', ax=ax_avg, title='Panier moyen par heure')
        plt.tight_layout(); st.pyplot(fig_temporal_avg)

        # Semaine vs week-end
        st.subheader("Semaine vs Week-end")
        df['WEEKEND'] = df['DATETIME'].dt.dayofweek >= 5
        wk = df['WEEKEND'].value_counts().rename({False:'Semaine',True:'Week-end'})
        fig_wk, ax_wk = plt.subplots()
        wk.plot(kind='bar', ax=ax_wk); plt.tight_layout(); st.pyplot(fig_wk)

                # 3. Analyse spatiale par code postal
        st.header("3. Analyse spatiale par code postal")
        spatial = df.groupby('CODE POSTAL').agg(
            tx_count=('MONTANT','size'),
            ca=('MONTANT','sum')
        ).reset_index()
        spatial['panier_moy'] = spatial['ca'] / spatial['tx_count']
        spatial = spatial.merge(df_geo, on='CODE POSTAL', how='left')
        st.dataframe(spatial)
        map_data = spatial.dropna(subset=['latitude','longitude'])
        # Capture matplotlib scatter pour PPT
        fig_spatial, ax_spat = plt.subplots()
        spatial.plot.scatter(x='longitude', y='latitude', s=spatial['panier_moy']/spatial['panier_moy'].max()*200, ax=ax_spat)
        ax_spat.set_title('Carte spatiale (taille ~ panier moyen)')
        plt.tight_layout()
        st.pyplot(fig_spatial)
        if not map_data.empty:
            st.subheader("Heatmap panier moyen")
            deck1 = pdk.Deck(
                map_style='mapbox://styles/mapbox/light-v10',
                initial_view_state=pdk.ViewState(latitude=46.5, longitude=2.5, zoom=5),
                layers=[pdk.Layer(
                    'HeatmapLayer',
                    data=map_data,
                    get_position=['longitude','latitude'],
                    get_weight='panier_moy',
                    radiusPixels=50
                )]
            )
            st.pydeck_chart(deck1)
            st.subheader("Scatter : montant total")
            deck2 = pdk.Deck(
                map_style='mapbox://styles/mapbox/light-v10',
                initial_view_state=pdk.ViewState(latitude=46.5, longitude=2.5, zoom=5),
                layers=[pdk.Layer(
                    'ScatterplotLayer',
                    data=map_data,
                    get_position=['longitude','latitude'],
                    get_radius='tx_count * 500',
                    get_fill_color=[200,30,0,160],
                    pickable=True
                )]
            )
            st.pydeck_chart(deck2)

        # 4. Diagnostics météo Diagnostics
        st.header("4. Diagnostics météo")
        corr = df[['TEMP','MONTANT']].corr().loc['TEMP','MONTANT']
        st.write(f"Corrélation temp/montant: {corr:.2f}")
        fig_corr, ax_corr = plt.subplots()
        ax_corr.scatter(df['TEMP'],df['MONTANT'],alpha=0.3)
        if len(df[['TEMP','MONTANT']].dropna())>1:
            coef=np.polyfit(df['TEMP'],df['MONTANT'],1)
            ax_corr.plot(df['TEMP'],coef[0]*df['TEMP']+coef[1],color='red')
        plt.tight_layout(); st.pyplot(fig_corr)

        # 5. Segmentation
        st.header("5. Segmentation client")
        feats = df[['MONTANT','HOUR','TEMP']].dropna()
        fig_seg=None
        if len(feats)>=3:
            X=StandardScaler().fit_transform(feats)
            labels=KMeans(n_clusters=3,random_state=42).fit_predict(X)
            df['cluster']=labels
            fig_seg, ax_seg = plt.subplots()
            pd.Series(labels).value_counts().sort_index().plot(kind='bar',ax=ax_seg,
                title='Clusters KMeans'); plt.tight_layout(); st.pyplot(fig_seg)

        st.info("Sections prédictives à venir.")

        # Vérification template
        if file_template is None:
            st.error("Upload du modèle PPTX requis."); st.stop()
        tmp_tpl=tempfile.NamedTemporaryFile(suffix='.pptx',delete=False)
        tmp_tpl.write(file_template.read()); tmp_tpl.flush()
        prs=Presentation(tmp_tpl.name)
        cover=prs.slides[0]
        cover.shapes.title.text="Smile Magic Report"
        try: cover.placeholders[1].text=pd.Timestamp.today().strftime("%d %B %Y")
        except: pass

                # Map sections à figures
        sections = {
            "Indicateurs descriptifs": fig_desc,
            "Temporalité - volume": fig_temporal_nb,
            "Temporalité - panier": fig_temporal_avg,
            "Analyse spatiale": fig_spatial,
            "Diagnostics météo": fig_corr,
            "Segmentation client": fig_seg
        }

        for title,fig in sections.items():
            slide=prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text=title
            prompt=f"Explique brièvement les résultats de la section '{title}'."
            resp=openai.chat.completions.create(model="gpt-4",
                messages=[{"role":"system","content":"Tu es un expert BI."},{"role":"user","content":prompt}]
            )
            exp=resp.choices[0].message.content
            try: slide.placeholders[1].text=exp
            except:
                tb=slide.shapes.add_textbox(Inches(1),Inches(1.5),Inches(8),Inches(1.5))
                tf=tb.text_frame; tf.text=exp
                for p in tf.paragraphs:
                    for run in p.runs: run.font.size=Pt(12)
            if fig:
                tmp= tempfile.NamedTemporaryFile(suffix='.png',delete=False)
                fig.savefig(tmp.name,bbox_inches='tight')
                slide.shapes.add_picture(tmp.name,Inches(1),Inches(2),width=Inches(8))

        out='smile_report_bi.pptx'; prs.save(out)
        st.success(f"🎉 Présentation générée : {out}")
        with open(out,'rb') as f: data=f.read()
        st.download_button("⬇️ Télécharger le rapport",data=data,file_name=out,mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # --- Fonctions LLM ---
    def get_mean_by_type(): return {'mean_by_type':df.groupby('TYPE_COMMERCE')['MONTANT'].mean().to_dict()}
    def get_top_merchants(by='transactions',top_n=5):
        data=df['MARCHAND'].value_counts().head(top_n).to_dict() if by=='transactions' else df.groupby('MARCHAND')['MONTANT'].sum().nlargest(top_n).to_dict()
        return {'top_merchants':data}
    def get_correlation():return{'corr':float(df[['TEMP','MONTANT']].corr().loc['TEMP','MONTANT'])}
    def get_distribution(percentiles=[0.1,0.25,0.5,0.75,0.9]):
        d=df['MONTANT'].describe(percentiles=percentiles)
        res={str(p):float(d[f'{int(p*100)}%'])for p in percentiles};res['mean']=float(d['mean']);return{'dist':res}
    def get_tpe_count_by_type():return{'records':df.groupby(['TYPE_COMMERCE','MODELE_TERMINAL']).size().reset_index(name='count').to_dict(orient='records')}

    functions=[
        {'name':'get_mean_by_type','description':'Panier moyen par type','parameters':{'type':'object','properties':{},'required':[]}},
        {'name':'get_top_merchants','description':'Top marchands','parameters':{'type':'object','properties':{'by':{'type':'string','enum':['transactions','revenue']},'top_n':{'type':'integer'}},'required':[]}},
        {'name':'get_correlation','description':'Corrélation temp','parameters':{'type':'object','properties':{},'required':[]}},
        {'name':'get_distribution','description':'Distribution percentiles','parameters':{'type':'object','properties':{'percentiles':{'type':'array','items':{'type':'number'}}},'required':[]}},
        {'name':'get_tpe_count_by_type','description':'TPE count','parameters':{'type':'object','properties':{},'required':[]}}
    ]

    st.header("💬 Interrogez l'agent BI")
    if 'chat_history' not in st.session_state or not isinstance(st.session_state.chat_history,list):st.session_state.chat_history=[]
    user_input=st.text_input("Votre question :",key='chat_input')
    if user_input and st.button("Envoyer",key='send_btn'):
        st.session_state.chat_history.append({'role':'user','content':user_input})
        messages=[{'role':'system','content':'Tu es un assistant BI expert.'}]+st.session_state.chat_history
        resp=openai.chat.completions.create(model='gpt-4-0613',messages=messages,functions=functions,function_call='auto')
        msg=resp.choices[0].message
        if getattr(msg,'function_call',None):fn=msg.function_call.name;args=json.loads(msg.function_call.arguments or '{}');res=globals()[fn](**args);st.session_state.chat_history.append({'role':'function','name':fn,'content':json.dumps(res)});follow=openai.chat.completions.create(model='gpt-4-0613',messages=[*messages,{'role':'assistant','content':None,'function_call':msg.function_call},{'role':'function','name':fn,'content':json.dumps(res)}]);reply=follow.choices[0].message.content
        else: reply=msg.content
        st.session_state.chat_history.append({'role':'assistant','content':reply})
    for chat in st.session_state.chat_history:
        if chat['role'] in ['user','assistant']:icon='👤' if chat['role']=='user' else '🤖';st.markdown(f"**{icon}** {chat['content']}")
else:
    st.warning("Chargez d'abord les 3 fichiers Excel.")

# === 5. Smile Vision 🪄 — Agent de prévision ===
st.header("✨ Smile Vision 🪄 — Prédiction transactions 7 prochains jours")

# Initialisation de l’historique de vision (pour conservation de la question)
if "vision_history" not in st.session_state:
    st.session_state.vision_history = []

vision_input = st.text_input("Pose ta question à Smile Vision :", key="vision_input")
if st.button("🔮 Envoyer à Smile Vision") and vision_input:
    st.session_state.vision_history.append(vision_input)
    # On attend forcément une phrase de type « prédit les transactions… »
    st.info("🔍 L'agent Smile Vision réfléchit…")

    # 1) Préparation de l’historique journalier par type de commerce
    # On part du df complet (transactions + type + météo journalière)
    # Recalculez la météo journalière moyenne si besoin :
    df_weather_full = read_file(file_weather)  # Excel complet, avec colonnes 'Date', 'CODE POSTAL', 'Température'
    # Normalisation des noms
    df_weather_full.rename(columns={
        'Date': 'DATE', 'Température': 'TEMP', 'CODE_POSTAL': 'CODE POSTAL'
    }, inplace=True)
    df_weather_full['DATE'] = pd.to_datetime(df_weather_full['DATE'], dayfirst=True, errors='coerce')
    
    # On reconstitue un df complet journalier
    df_tx_j = df_tx.copy()
    df_tx_j['DATE'] = df_tx_j['DATETIME'].dt.date
    df_tx_j = df_tx_j.merge(
        df_merch.rename(columns={'Organization_type':'TYPE_COMMERCE'})[['REF_MARCHAND','TYPE_COMMERCE']],
        on='REF_MARCHAND', how='left'
    ).merge(
        df_weather_full[['DATE','CODE POSTAL','TEMP']],
        on=['DATE','CODE POSTAL'], how='left'
    )

    hist = (
        df_tx_j
        .groupby(['DATE','TYPE_COMMERCE'])
        .agg(tx_count=('MONTANT','size'),
             temp_moy=('TEMP','mean'))
        .reset_index()
    )

    # 2) Entraînement d’un modèle linéaire par type de commerce
    from sklearn.linear_model import LinearRegression
    models = {}
    for typ in hist['TYPE_COMMERCE'].unique():
        sub = hist[hist['TYPE_COMMERCE']==typ].dropna()
        if len(sub) >= 2:
            lr = LinearRegression().fit(sub[['temp_moy']], sub['tx_count'])
            models[typ] = lr

    # 3) Constitution des 7 jours à venir + météo
    # Pour POC, on prend df_weather_full jusqu’à 7 jours après la dernière date connue
    last_date = hist['DATE'].max()
    futura = df_weather_full[df_weather_full['DATE'] > last_date] \
             .sort_values('DATE') \
             .drop_duplicates('DATE') \
             .head(7) \
             .loc[:, ['DATE','TEMP']]

    # 4) Calcul des prévisions
    preds = []
    for typ, lr in models.items():
        y_pred = lr.predict(futura[['TEMP']])
        tmp = futura.copy()
        tmp['TYPE_COMMERCE'] = typ
        tmp['tx_pred'] = y_pred
        preds.append(tmp)
    df_pred = pd.concat(preds, ignore_index=True)

    # 5) Affichage du graphique historique vs prévision
    fig, ax = plt.subplots(figsize=(8,4))
    for typ in hist['TYPE_COMMERCE'].unique():
        sub_h = hist[hist['TYPE_COMMERCE']==typ]
        ax.plot(sub_h['DATE'], sub_h['tx_count'], label=f"{typ} (historique)")
        sub_p = df_pred[df_pred['TYPE_COMMERCE']==typ]
        ax.plot(sub_p['DATE'], sub_p['tx_pred'], '--', label=f"{typ} (prévu)")
    ax.set_xlabel("Date")
    ax.set_ylabel("Nb transactions")
    ax.set_title("Historique vs prévisions 7 jours par type de commerce")
    ax.legend(loc='upper left', fontsize='small')
    plt.xticks(rotation=45)
    plt.tight_layout()
    st.pyplot(fig)
